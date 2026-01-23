import requests
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from .schemas import GeminiOutline

def get_image(keyword):
    if not keyword or len(keyword.strip()) < 2:
        return None
    try:
        url = f"https://source.unsplash.com/800x600/?{keyword}"
        resp = requests.get(url, timeout=2)
        if resp.status_code == 200 and len(resp.content) > 1000:
            return BytesIO(resp.content)
    except:
        pass
    return None

def fetch_all_images(keywords):
    """Fetch all images in parallel - much faster!"""
    images = {}
    with ThreadPoolExecutor(max_workers=5) as executor:
        future_to_kw = {executor.submit(get_image, kw): kw for kw in keywords}
        for future in as_completed(future_to_kw):
            kw = future_to_kw[future]
            try:
                images[kw] = future.result()
            except:
                images[kw] = None
    return images

def build_pptx(outline: GeminiOutline) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Fetch all images in parallel FIRST
    keywords = [s.image_keyword for s in outline.slides if s.image_keyword]
    images = fetch_all_images(keywords)

    # === Title Slide ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 40, 80)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = outline.presentation_title
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if outline.subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
        tf2 = sub_box.text_frame
        tf2.paragraphs[0].text = outline.subtitle
        tf2.paragraphs[0].font.size = Pt(24)
        tf2.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # === Content Slides ===
    colors = [
        RGBColor(0, 102, 153),
        RGBColor(0, 128, 128),
        RGBColor(70, 130, 180),
        RGBColor(100, 100, 160),
        RGBColor(80, 120, 140),
    ]

    for idx, s in enumerate(outline.slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        color = colors[idx % len(colors)]

        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()

        # Slide title
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(10), Inches(0.7))
        tf = title_box.text_frame
        tf.paragraphs[0].text = s.title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12.3), Inches(0.25), Inches(0.8), Inches(0.7))
        tf_num = num_box.text_frame
        tf_num.paragraphs[0].text = str(idx + 1)
        tf_num.paragraphs[0].font.size = Pt(20)
        tf_num.paragraphs[0].font.bold = True
        tf_num.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf_num.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Try to add image
        has_image = False
        img_data = images.get(s.image_keyword)
        if img_data:
            try:
                img_data.seek(0)
                slide.shapes.add_picture(img_data, Inches(8.3), Inches(1.4), width=Inches(4.5))
                has_image = True
            except:
                pass

        # Bullets
        bullet_width = Inches(7.3) if has_image else Inches(12.3)
        bullet_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.4), bullet_width, Inches(5.5))
        tf_bullets = bullet_box.text_frame
        tf_bullets.word_wrap = True

        for i, bullet in enumerate(s.bullets):
            if i == 0:
                p = tf_bullets.paragraphs[0]
            else:
                p = tf_bullets.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(40, 40, 40)
            p.space_after = Pt(10)

        # Bottom accent
        bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3))
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = color
        bottom.line.fill.background()

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()